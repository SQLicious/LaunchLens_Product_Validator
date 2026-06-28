"""Generate the LaunchLens submission deck (slides/LaunchLens.pptx).

Aesthetic emulates the course brief video: deep-navy plates with soft corner
glows (slides/assets/bg, built by build_backgrounds.py), Poppins type, amber +
emerald + lavender accents, rounded cards with subtle borders, pill badges and
big ghost watermark words. Text stays as real PowerPoint runs (editable).

Run order:
    python slides/build_backgrounds.py   # plates
    python slides/build_deck.py          # this deck
Export to PDF from PowerPoint/Google Slides for submission.
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(__file__)
BG = os.path.join(HERE, "assets", "bg")
ASSETS = os.path.join(HERE, "assets")

# ---- palette ----
WHITE = RGBColor(0xF5, 0xF7, 0xFB)
AMBER = RGBColor(0xF5, 0xA6, 0x23)
EMERALD = RGBColor(0x34, 0xD3, 0x99)
LAV = RGBColor(0xA5, 0xB4, 0xFC)
KICKER = RGBColor(0x8A, 0x92, 0xA8)
BODY = RGBColor(0x9A, 0xA3, 0xB2)
DIM = RGBColor(0x6B, 0x74, 0x88)
CARD_FILL = RGBColor(0x12, 0x18, 0x26)
CARD_LINE = RGBColor(0x27, 0x30, 0x42)
WM = RGBColor(0x18, 0x20, 0x30)

# ---- fonts (installed per-user) ----
F_XB = "Poppins ExtraBold"
F_SB = "Poppins SemiBold"
F_MD = "Poppins Medium"
F_RG = "Poppins"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------- helpers
def slide(plate="bg_content.png"):
    s = prs.slides.add_slide(BLANK)
    s.shapes.add_picture(os.path.join(BG, plate), 0, 0, SW, SH)
    return s


def _track(run, pts):
    run._r.get_or_add_rPr().set("spc", str(int(pts * 100)))


def _run(p, text, size, color, font=F_RG, spc=0.0):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.name = font
    if spc:
        _track(r, spc)
    return r


def box(s, x, y, w, h, anchor=None):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def watermark(s, text, size, x, y, color=WM):
    tf = box(s, x, y, Inches(12), Inches(3.5))
    _run(tf.paragraphs[0], text, size, color, font=F_XB)


def kicker(s, text, x=Inches(0.9), y=Inches(0.95), color=KICKER):
    tf = box(s, x, y, Inches(11.5), Inches(0.45))
    _run(tf.paragraphs[0], text.upper(), 14, color, font=F_SB, spc=3.2)


def headline(s, segments, x=Inches(0.88), y=Inches(1.35), size=40, width=11.6):
    """segments: list of (text, color). Rendered as one wrapping paragraph."""
    tf = box(s, x, y, Inches(width), Inches(2.2))
    p = tf.paragraphs[0]
    p.line_spacing = 1.04
    for text, color in segments:
        _run(p, text, size, color, font=F_XB)
    return tf


def underline(s, x, y, w=Inches(1.25), color=AMBER, h=Pt(4)):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp


def card(s, x, y, w, h, fill=CARD_FILL, line=CARD_LINE, radius=0.07):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = radius
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def hbar(s, x, y, w, color=CARD_LINE, h=Pt(1.5)):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp


def pill(s, x, y, lead, strong, w=Inches(3.1), h=Inches(0.62), icon=None):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.5
    shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor(0x14, 0x24, 0x1E)
    shp.line.color.rgb = RGBColor(0x1F, 0x3A, 0x30); shp.line.width = Pt(1)
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = False
    tf.margin_left = Inches(0.62 if icon else 0.25); tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if icon else PP_ALIGN.CENTER
    _run(p, lead + " ", 16, BODY, font=F_MD)
    _run(p, strong, 16, WHITE, font=F_SB)
    if icon:
        try:
            s.shapes.add_picture(icon, Emu(int(x) + int(Inches(0.22))),
                                 Emu(int(y) + int((int(h) - int(Inches(0.32))) / 2)),
                                 height=Inches(0.32))
        except Exception:
            pass
    return shp


# =================================================================== SLIDES

# 1 -- COVER ---------------------------------------------------------------
s = slide("bg_title.png")
watermark(s, "LAUNCH", 200, Inches(1.0), Inches(4.55))
kicker(s, "Assignment 03 · Startup Build", y=Inches(2.35))
tf = box(s, Inches(0.9), Inches(2.85), Inches(11), Inches(0.6))
_run(tf.paragraphs[0], "Before every launch, one scary question.", 24, BODY, font=F_MD)
headline(s, [("Will anyone actually ", WHITE), ("buy", AMBER), (" this?", WHITE)],
         y=Inches(3.45), size=58, width=11.4)
underline(s, Inches(0.92), Inches(5.95), w=Inches(1.6))
tf = box(s, Inches(0.92), Inches(6.55), Inches(11), Inches(0.5))
_run(tf.paragraphs[0], "LaunchLens  ·  Ruby Gunna  ·  LangGraph + SerpApi + Oxylabs", 14, DIM, font=F_MD)

# 2 -- MEET THE PRODUCT ----------------------------------------------------
s = slide("bg_content.png")
kicker(s, "The product")
headline(s, [("LaunchLens turns that fear into a ", WHITE), ("verdict.", EMERALD)],
         y=Inches(1.5), size=40)
underline(s, Inches(0.92), Inches(3.0), color=EMERALD)
tf = box(s, Inches(0.9), Inches(3.45), Inches(11.4), Inches(2.0))
p = tf.paragraphs[0]; p.line_spacing = 1.3
_run(p, "A market-intelligence copilot for founders. Ask in plain English, and it "
        "fuses live demand signals from Google with live supply signals from Amazon "
        "into one ", 21, BODY, font=F_RG)
_run(p, "Go / No-Go / Niche", 21, WHITE, font=F_SB)
_run(p, " call — with a price band, a positioning angle, and a confidence level.", 21, BODY, font=F_RG)

# 3 -- PROBLEM -------------------------------------------------------------
s = slide("bg_content.png")
kicker(s, "The problem")
headline(s, [("The signals never meet.", WHITE)], y=Inches(1.4), size=38)
underline(s, Inches(0.92), Inches(2.55))
probs = [
    ("Demand lives on Google", "Trends, News, Shopping — is anyone searching, is it seasonal, what does it cost?", AMBER),
    ("Supply lives on Amazon", "Who already sells it, at what price, and what do the reviews complain about?", EMERALD),
    ("Nobody fuses them", "A trend chart says 'maybe'. A competitor list says 'crowded'. The answer is in the overlap.", LAV),
]
cy = Inches(3.05); cw = Inches(3.74); gap = Inches(0.2); cx = Inches(0.88); chh = Inches(2.7)
for i, (title, desc, col) in enumerate(probs):
    x = Emu(int(cx) + i * (int(cw) + int(gap)))
    card(s, x, cy, cw, chh)
    hbar(s, Emu(int(x) + Inches(0.28)), Emu(int(cy) + Inches(0.32)), Inches(0.5), color=col, h=Pt(3))
    tf = box(s, Emu(int(x) + Inches(0.28)), Emu(int(cy) + Inches(0.55)), Inches(3.2), Inches(1.9))
    _run(tf.paragraphs[0], title, 19, WHITE, font=F_SB)
    p = tf.add_paragraph(); p.space_before = Pt(10); p.line_spacing = 1.25
    _run(p, desc, 14, BODY, font=F_RG)

# 4 -- WHAT IT DOES --------------------------------------------------------
s = slide("bg_content.png")
kicker(s, "What it does")
headline(s, [("One question in, one judgement out.", WHITE)], y=Inches(1.4), size=36)
underline(s, Inches(0.92), Inches(2.5), color=EMERALD)
steps = [
    ("01", "Ask", "“Should I launch a $20 meal-prep set in the US?”", AMBER),
    ("02", "Gather", "Pulls demand + supply live, in parallel.", EMERALD),
    ("03", "Judge", "Go / No-Go / Niche + price band + confidence.", LAV),
    ("04", "Remember", "Follow-ups keep context across the chat.", AMBER),
]
cy = Inches(3.0); cw = Inches(2.82); gap = Inches(0.18); cx = Inches(0.88); chh = Inches(2.7)
for i, (num, title, desc, col) in enumerate(steps):
    x = Emu(int(cx) + i * (int(cw) + int(gap)))
    card(s, x, cy, cw, chh)
    tf = box(s, Emu(int(x) + Inches(0.26)), Emu(int(cy) + Inches(0.28)), Inches(2.3), Inches(2.2))
    _run(tf.paragraphs[0], num, 22, col, font=F_XB)
    p = tf.add_paragraph(); p.space_before = Pt(6); _run(p, title, 18, WHITE, font=F_SB)
    p = tf.add_paragraph(); p.space_before = Pt(8); p.line_spacing = 1.25
    _run(p, desc, 13.5, BODY, font=F_RG)

# 5 -- DEMO FLOW -----------------------------------------------------------
s = slide("bg_content.png")
kicker(s, "Demo flow")
headline(s, [("A founder conversation, end to end.", WHITE)], y=Inches(1.4), size=34)
underline(s, Inches(0.92), Inches(2.5))
flow = [
    ("“…worth launching a 32oz steel bottle under $40?”", "full fan-out → verdict", AMBER),
    ("“What are reviewers complaining about?”", "mined Amazon complaints", EMERALD),
    ("“What about a premium glass version at $35?”", "remembers the product — memory", LAV),
    ("“What did we conclude overall?”", "summarized recall", AMBER),
]
ry = Inches(2.95); rh = Inches(0.92); rgap = Inches(0.16); rx = Inches(0.88); rw = Inches(11.55)
for i, (q, tag, col) in enumerate(flow):
    y = Emu(int(ry) + i * (int(rh) + int(rgap)))
    card(s, rx, y, rw, rh, radius=0.18)
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(int(rx) + Inches(0.16)),
                             Emu(int(y) + Inches(0.16)), Inches(0.08), Inches(0.6))
    bar.adjustments[0] = 0.5; bar.fill.solid(); bar.fill.fore_color.rgb = col
    bar.line.fill.background(); bar.shadow.inherit = False
    tf = box(s, Emu(int(rx) + Inches(0.5)), y, Inches(7.3), rh, anchor=MSO_ANCHOR.MIDDLE)
    _run(tf.paragraphs[0], q, 17, WHITE, font=F_MD)
    tf = box(s, Emu(int(rx) + Inches(7.9)), y, Inches(3.4), rh, anchor=MSO_ANCHOR.MIDDLE)
    pr = tf.paragraphs[0]; pr.alignment = PP_ALIGN.RIGHT
    _run(pr, tag, 14, col, font=F_SB)

# 6 -- ARCHITECTURE (matches brief "wired into one graph") -----------------
s = slide("bg_content.png")
pill(s, Inches(0.88), Inches(0.85), "Built on", "LangGraph", w=Inches(3.25),
     icon=os.path.join(ASSETS, "langgraph.png"))
headline(s, [("Everything, wired into one graph.", WHITE)], y=Inches(1.85), size=38)
arch = [
    ("Router", "read the question", AMBER),
    ("Fan-out", "pull sources in parallel", EMERALD),
    ("Agent + Tools", "SerpApi & Oxylabs", AMBER),
    ("Memory", "short-term + summarize", EMERALD),
]
cy = Inches(3.7); cw = Inches(2.72); gap = Inches(0.27); cx = Inches(0.88); chh = Inches(1.7)
centers_y = int(cy) + int(chh) // 2
for i, (title, desc, col) in enumerate(arch):
    x = int(cx) + i * (int(cw) + int(gap))
    if i > 0:
        hbar(s, Emu(x - int(gap) - Inches(0.02)), Emu(centers_y), Emu(int(gap) + Inches(0.04)), color=CARD_LINE)
    card(s, Emu(x), cy, cw, chh)
    tf = box(s, Emu(x + int(Inches(0.3))), Emu(int(cy) + int(Inches(0.34))), Inches(2.2), Inches(1.1))
    _run(tf.paragraphs[0], title, 19, col, font=F_SB)
    p = tf.add_paragraph(); p.space_before = Pt(8); _run(p, desc, 13.5, BODY, font=F_RG)

# 7 -- FIVE CONCEPTS -------------------------------------------------------
s = slide("bg_content.png")
kicker(s, "Under the hood")
headline(s, [("The five concepts, mapped to code.", WHITE)], y=Inches(1.4), size=34)
underline(s, Inches(0.92), Inches(2.5), color=EMERALD)
concepts = [
    ("1", "Graph & state", "Typed StateGraph + reducer channels", "state.py 16/28 · graph.py 28", AMBER),
    ("2", "Fan-out (parallel)", "Router returns a 3-node list, results merge", "graph.py 49–65 · nodes.py 98/104/110", EMERALD),
    ("3", "Routing", "Intent → conditional edges, chat = default", "router.py 65/81 · graph.py 49", LAV),
    ("4", "Agent + tools", "LLM ⇄ ToolNode loop over 6 tools", "nodes.py 135 · graph.py 41/68", AMBER),
    ("5", "Short-term memory", "Checkpointer + summarization", "main.py 140 · nodes.py 64", EMERALD),
]
ry = Inches(2.9); rh = Inches(0.78); rgap = Inches(0.12); rx = Inches(0.88); rw = Inches(11.55)
for i, (num, title, desc, ref, col) in enumerate(concepts):
    y = Emu(int(ry) + i * (int(rh) + int(rgap)))
    card(s, rx, y, rw, rh, radius=0.16)
    tf = box(s, Emu(int(rx) + Inches(0.3)), y, Inches(0.7), rh, anchor=MSO_ANCHOR.MIDDLE)
    _run(tf.paragraphs[0], num, 22, col, font=F_XB)
    tf = box(s, Emu(int(rx) + Inches(1.05)), y, Inches(3.3), rh, anchor=MSO_ANCHOR.MIDDLE)
    _run(tf.paragraphs[0], title, 17, WHITE, font=F_SB)
    tf = box(s, Emu(int(rx) + Inches(4.3)), y, Inches(4.0), rh, anchor=MSO_ANCHOR.MIDDLE)
    _run(tf.paragraphs[0], desc, 14, BODY, font=F_RG)
    tf = box(s, Emu(int(rx) + Inches(8.3)), y, Inches(3.0), rh, anchor=MSO_ANCHOR.MIDDLE)
    pr = tf.paragraphs[0]; pr.alignment = PP_ALIGN.RIGHT
    _run(pr, ref, 11.5, DIM, font=F_MD)

# 8 -- DATA FUSION (core, matches brief GO frame) --------------------------
s = slide("bg_core.png")
watermark(s, "GO", 230, Inches(0.4), Inches(0.1))
kicker(s, "Demand, market and price — fused into one answer", y=Inches(1.5))
rows = [
    ("DEMAND", "is the interest really there?", AMBER),
    ("MARKET", "who already wins, and why?", EMERALD),
    ("PRICE", "where would yours land?", LAV),
]
ry = Inches(2.25); rh = Inches(0.86); rgap = Inches(0.2); rx = Inches(0.95); rw = Inches(10.6)
for i, (label, q, col) in enumerate(rows):
    y = Emu(int(ry) + i * (int(rh) + int(rgap)))
    card(s, rx, y, rw, rh, radius=0.14)
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, Emu(int(y) + Inches(0.13)),
                             Inches(0.09), Inches(0.6))
    bar.adjustments[0] = 0.5; bar.fill.solid(); bar.fill.fore_color.rgb = col
    bar.line.fill.background(); bar.shadow.inherit = False
    tf = box(s, Emu(int(rx) + Inches(0.4)), y, Inches(2.2), rh, anchor=MSO_ANCHOR.MIDDLE)
    _run(tf.paragraphs[0], label, 15, col, font=F_SB, spc=2.5)
    tf = box(s, Emu(int(rx) + Inches(2.7)), y, Inches(7.5), rh, anchor=MSO_ANCHOR.MIDDLE)
    _run(tf.paragraphs[0], q, 20, WHITE, font=F_SB)
# big GO + caption
tf = box(s, Inches(0.9), Inches(5.55), Inches(3.0), Inches(1.4))
_run(tf.paragraphs[0], "GO", 80, EMERALD, font=F_XB)
tf = box(s, Inches(2.95), Inches(5.55), Inches(8.0), Inches(1.4), anchor=MSO_ANCHOR.MIDDLE)
_run(tf.paragraphs[0], "or no-go. One clear call.", 26, BODY, font=F_RG)

# 9 -- TOOLS ---------------------------------------------------------------
s = slide("bg_content.png")
kicker(s, "Six tools, two providers")
headline(s, [("Slim JSON in, evidence out.", WHITE)], y=Inches(1.4), size=34)
underline(s, Inches(0.92), Inches(2.5))
cols = [
    ("DEMAND · SerpApi", AMBER, [
        ("google_trends", "interest + seasonality, slope-classified"),
        ("google_news", "launches, recalls, competitor moves"),
        ("google_shopping", "cross-retailer price band"),
    ]),
    ("SUPPLY · Oxylabs", EMERALD, [
        ("amazon_search", "live listings, prices, review counts"),
        ("amazon_product", "per-ASIN mined review complaints"),
        ("amazon_bestsellers", "how entrenched the competition is"),
    ]),
]
cw = Inches(5.6); gap = Inches(0.35); cx = Inches(0.88); cy = Inches(2.95); chh = Inches(3.7)
for i, (head, col, items) in enumerate(cols):
    x = Emu(int(cx) + i * (int(cw) + int(gap)))
    card(s, x, cy, cw, chh)
    tf = box(s, Emu(int(x) + Inches(0.35)), Emu(int(cy) + Inches(0.3)), Inches(5.0), Inches(0.5))
    _run(tf.paragraphs[0], head, 15, col, font=F_SB, spc=1.5)
    tf = box(s, Emu(int(x) + Inches(0.35)), Emu(int(cy) + Inches(0.95)), Inches(5.0), Inches(2.6))
    for j, (name, desc) in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        _run(p, name, 16, WHITE, font=F_SB)
        p2 = tf.add_paragraph(); p2.space_after = Pt(10)
        _run(p2, desc, 13.5, BODY, font=F_RG)
tf = box(s, Inches(0.9), Inches(6.95), Inches(11.5), Inches(0.4))
_run(tf.paragraphs[0], "Every tool returns a few fields, never the raw payload — tokens + context stay bounded.",
     13, DIM, font=F_MD)

# 10 -- MEMORY -------------------------------------------------------------
s = slide("bg_content.png")
kicker(s, "Short-term memory")
headline(s, [("Survives restarts and long chats.", WHITE)], y=Inches(1.4), size=34)
underline(s, Inches(0.92), Inches(2.5), color=EMERALD)
mem = [
    ("Checkpointer", "SqliteSaver persists full state after every node, keyed by thread_id — quit and resume.", AMBER),
    ("Summarization", "Old plain messages fold into a rolling summary once the chat grows.", EMERALD),
    ("History intact", "Only Human/AI text is compressed; tool-call sequences are never broken.", LAV),
    ("Flat cost", "Per-turn cost stays roughly constant instead of growing with the conversation.", AMBER),
]
cy = Inches(3.0); cw = Inches(5.6); gap = Inches(0.35); chh = Inches(1.65)
for i, (title, desc, col) in enumerate(mem):
    x = Emu(int(Inches(0.88)) + (i % 2) * (int(cw) + int(gap)))
    y = Emu(int(cy) + (i // 2) * (int(chh) + int(Inches(0.25))))
    card(s, x, y, cw, chh)
    hbar(s, Emu(int(x) + Inches(0.32)), Emu(int(y) + Inches(0.3)), Inches(0.5), color=col, h=Pt(3))
    tf = box(s, Emu(int(x) + Inches(0.32)), Emu(int(y) + Inches(0.52)), Inches(5.0), Inches(1.0))
    _run(tf.paragraphs[0], title, 18, WHITE, font=F_SB)
    p = tf.add_paragraph(); p.space_before = Pt(8); p.line_spacing = 1.2
    _run(p, desc, 13.5, BODY, font=F_RG)

# 11 -- VERDICT + CONFIDENCE ----------------------------------------------
s = slide("bg_content.png")
kicker(s, "Verdict + confidence")
headline(s, [("A call you can trust — and challenge.", WHITE)], y=Inches(1.4), size=34)
underline(s, Inches(0.92), Inches(2.5))
ver = [
    ("Conditional, not absolute", "“Go IF COGS < $X and a defensible feature exists.”", AMBER),
    ("Confidence level", "High / Medium / Low — lowered when a source comes back empty.", EMERALD),
    ("Trend sanity", "Treats Trends as a relative, seasonal index — 5→10 isn't 'doubling'.", LAV),
    ("Unit-economics aware", "Prompts margin reasoning: COGS, fees, ad cost before a Go.", AMBER),
]
ry = Inches(2.95); rh = Inches(0.92); rgap = Inches(0.16); rx = Inches(0.88); rw = Inches(11.55)
for i, (title, desc, col) in enumerate(ver):
    y = Emu(int(ry) + i * (int(rh) + int(rgap)))
    card(s, rx, y, rw, rh, radius=0.16)
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(int(rx) + Inches(0.16)),
                             Emu(int(y) + Inches(0.16)), Inches(0.08), Inches(0.6))
    bar.adjustments[0] = 0.5; bar.fill.solid(); bar.fill.fore_color.rgb = col
    bar.line.fill.background(); bar.shadow.inherit = False
    tf = box(s, Emu(int(rx) + Inches(0.5)), y, Inches(3.6), rh, anchor=MSO_ANCHOR.MIDDLE)
    _run(tf.paragraphs[0], title, 17, WHITE, font=F_SB)
    tf = box(s, Emu(int(rx) + Inches(4.3)), y, Inches(7.0), rh, anchor=MSO_ANCHOR.MIDDLE)
    _run(tf.paragraphs[0], desc, 15, BODY, font=F_RG)

# 12 -- CLOSE --------------------------------------------------------------
s = slide("bg_core.png")
watermark(s, "SHIP", 200, Inches(0.7), Inches(4.4))
kicker(s, "Roadmap · thanks", color=EMERALD)
headline(s, [("From signal to ", WHITE), ("verdict", AMBER), (" — in one chat.", WHITE)],
         y=Inches(1.6), size=40)
underline(s, Inches(0.92), Inches(3.1), color=EMERALD)
tf = box(s, Inches(0.9), Inches(3.55), Inches(11.4), Inches(2.0))
for i, t in enumerate([
        "Next: a real FBA fee + COGS unit-economics module (margin, not just a price gap).",
        "Postgres checkpointer for multi-user prod · a small eval harness · richer streaming UI.",
        "Repo: github.com/SQLicious/LaunchLens_Product_Validator"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(12)
    _run(p, "→  ", 18, EMERALD, font=F_SB)
    _run(p, t, 18, BODY, font=F_RG)

out = os.path.join(HERE, "LaunchLens.pptx")
prs.save(out)
print("wrote", out, "—", len(prs.slides._sldIdLst), "slides")
