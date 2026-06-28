"""One-shot patch: insert a 'reasoning engine' (LLM/Qwen) slide into the
EXISTING LaunchLens.pptx without regenerating the deck (preserves manual edits).
Run once: `python slides/add_llm_slide.py`.
"""
import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt, Emu

HERE = os.path.dirname(__file__)
PPTX = os.path.join(HERE, "LaunchLens.pptx")
PLATE = os.path.join(HERE, "assets", "bg", "bg_content.png")

WHITE = RGBColor(0xF5, 0xF7, 0xFB); AMBER = RGBColor(0xF5, 0xA6, 0x23)
EMERALD = RGBColor(0x34, 0xD3, 0x99); LAV = RGBColor(0xA5, 0xB4, 0xFC)
KICKER = RGBColor(0x8A, 0x92, 0xA8); BODY = RGBColor(0x9A, 0xA3, 0xB2)
DIM = RGBColor(0x6B, 0x74, 0x88)
CARD_FILL = RGBColor(0x12, 0x18, 0x26); CARD_LINE = RGBColor(0x27, 0x30, 0x42)
F_XB = "Poppins ExtraBold"; F_SB = "Poppins SemiBold"; F_MD = "Poppins Medium"; F_RG = "Poppins"

prs = Presentation(PPTX)
SW, SH = prs.slide_width, prs.slide_height
s = prs.slides.add_slide(prs.slide_layouts[6])
s.shapes.add_picture(PLATE, 0, 0, SW, SH)


def box(x, y, w, h):
    tf = s.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def run(p, t, size, color, font=F_RG, spc=0.0):
    r = p.add_run(); r.text = t
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = font
    if spc:
        r._r.get_or_add_rPr().set("spc", str(int(spc * 100)))
    return r


# kicker + headline + underline
run(box(Inches(0.9), Inches(0.95), Inches(11.5), Inches(0.45)).paragraphs[0],
    "THE REASONING ENGINE", 14, KICKER, F_SB, spc=3.2)
run(box(Inches(0.88), Inches(1.4), Inches(11.6), Inches(1.0)).paragraphs[0],
    "Where the LLM thinks — and which one.", 34, WHITE, F_XB)
u = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.92), Inches(2.5), Inches(1.25), Pt(4))
u.fill.solid(); u.fill.fore_color.rgb = AMBER; u.line.fill.background(); u.shadow.inherit = False

cols = [
    ("WHERE IT RUNS", AMBER, [
        ("agent_node", "fuses demand + supply, picks tools, writes the verdict"),
        ("summarize_node", "compresses old turns into rolling memory"),
        ("router + verdict", "keyword + parser — no LLM, cheap & deterministic"),
    ]),
    ("MODEL & WHY", EMERALD, [
        ("qwen-plus · Qwen cloud (DashScope)", "OpenAI-compatible endpoint, temperature 0"),
        ("native tool / function calling", "required for the agent ⇄ tools loop"),
        ("low cost · big context · drop-in", "ChatOpenAI + base_url; swap via LLM_PROVIDER/MODEL"),
    ]),
]
cw = Inches(5.6); gap = Inches(0.35); cx = Inches(0.88); cy = Inches(2.95); chh = Inches(3.7)
for i, (head, col, items) in enumerate(cols):
    x = Emu(int(cx) + i * (int(cw) + int(gap)))
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, cy, cw, chh)
    card.adjustments[0] = 0.07
    card.fill.solid(); card.fill.fore_color.rgb = CARD_FILL
    card.line.color.rgb = CARD_LINE; card.line.width = Pt(1); card.shadow.inherit = False
    run(box(Emu(int(x) + int(Inches(0.35))), Emu(int(cy) + int(Inches(0.3))), Inches(5.0), Inches(0.5)).paragraphs[0],
        head, 15, col, F_SB, spc=1.5)
    tf = box(Emu(int(x) + int(Inches(0.35))), Emu(int(cy) + int(Inches(0.95))), Inches(5.0), Inches(2.6))
    for j, (name, desc) in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(4); run(p, name, 16, WHITE, F_SB)
        p2 = tf.add_paragraph(); p2.space_after = Pt(12); p2.line_spacing = 1.15
        run(p2, desc, 13.5, BODY, F_RG)

run(box(Inches(0.9), Inches(6.95), Inches(11.5), Inches(0.4)).paragraphs[0],
    "temperature 0 keeps Go / No-Go verdicts deterministic and reproducible.", 13, DIM, F_MD)

# move the new slide to second-to-last (just before the closing slide)
lst = prs.slides._sldIdLst
new = list(lst)[-1]
lst.remove(new)
lst.insert(len(list(lst)) - 1, new)

prs.save(PPTX)
print("inserted reasoning-engine slide; deck now", len(prs.slides._sldIdLst), "slides")
